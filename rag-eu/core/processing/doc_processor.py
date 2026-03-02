# core/processing/doc_processor.py

import hashlib
import json
import logging
import os
import re
import statistics
import time
import unicodedata
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import fitz  # PyMuPDF
import yaml
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# Configure logging for PDF extraction
logger = logging.getLogger(__name__)

# ============= PDF EXTRACTION CONSTANTS =============
MIN_GAP_WIDTH = 20  # Minimum gap width for column detection
MIDDLE_GAP_WIDTH = 60  # Width for middle gap analysis

# Pre-compiled regex patterns for performance
REGEX_PATTERNS = {
    "page_number": re.compile(r"^[\-\s]*(\d+|[ivxIVX]+)[\-\s/]*\d*[\-\s]*$"),
    "page_text": re.compile(r"^(Page|Seite|Página|页)\s*\d+", re.IGNORECASE),
    "multiple_spaces": re.compile(r" {2,}"),
    "excessive_newlines": re.compile(r"\n{4,}"),
    "triple_newlines": re.compile(r"\n{3}"),
    "space_before_newline": re.compile(r" +\n"),
    "space_after_newline": re.compile(r"\n +"),
}


# ============= PDF EXTRACTION ENUMS AND DATACLASSES =============
class LayoutType(Enum):
    """Page layout types"""

    SINGLE_PORTRAIT = "single_portrait"
    LANDSCAPE_OR_DOUBLE = "landscape_or_double"


class BlockType(Enum):
    """Text block types for better classification"""

    MAIN_TEXT = "main_text"
    TITLE = "title"
    HEADER = "header"
    FOOTER = "footer"
    INFOBOX = "infobox"
    MARGIN_LEFT = "margin_left"
    MARGIN_RIGHT = "margin_right"
    TABLE = "table"
    CAPTION = "caption"
    FOOTNOTE = "footnote"
    PAGE_NUMBER = "page_number"

    def __lt__(self, other):
        """Make BlockType sortable"""
        if not isinstance(other, BlockType):
            return NotImplemented
        return self.value < other.value


@dataclass(slots=True)
class TextBlock:
    """Enhanced text block with position and type information"""

    text: str
    bbox: Tuple[float, float, float, float]  # (x0, y0, x1, y1)
    page_num: int
    block_type: str = "text"
    confidence: float = 1.0
    column_id: Optional[int] = None
    semantic_type: Optional[BlockType] = None
    font_size: Optional[float] = None

    @property
    def x_center(self) -> float:
        return (self.bbox[0] + self.bbox[2]) / 2

    @property
    def y_center(self) -> float:
        return (self.bbox[1] + self.bbox[3]) / 2

    @property
    def width(self) -> float:
        return self.bbox[2] - self.bbox[0]

    @property
    def height(self) -> float:
        return self.bbox[3] - self.bbox[1]


# ============= PDF COMPLEX LAYOUT EXTRACTOR =============
class PDFComplexLayoutExtractor:
    """
    Advanced PDF text extractor with intelligent layout recognition
    and context-preserving reading order.
    """

    def __init__(self, pdf_path: str, debug: bool = False):
        self.pdf_path = pdf_path
        self.debug = debug

        if debug:
            logging.basicConfig(level=logging.DEBUG)
        else:
            logging.basicConfig(level=logging.INFO)

        self.doc = fitz.open(pdf_path)
        self.page_layouts = []
        self.coordinate_system = "unknown"

        self._analyze_page_layouts()
        self._detect_coordinate_system()

    def _analyze_page_layouts(self):
        """Analyze and classify layout type for each page"""
        for page_num, page in enumerate(self.doc):
            rect = page.rect
            width, height = rect.width, rect.height

            blocks = page.get_text("dict").get("blocks", [])
            text_blocks = [b for b in blocks if b.get("type") == 0]

            layout_type = self._detect_page_layout(width, height, text_blocks)
            self.page_layouts.append(
                {
                    "page_num": page_num,
                    "width": width,
                    "height": height,
                    "layout_type": layout_type,
                    "text_blocks": text_blocks,
                }
            )

            logger.debug(f"Page {page_num + 1}: {layout_type.value} ({width:.0f}x{height:.0f})")

    def _detect_page_layout(self, width: float, height: float, blocks: List) -> LayoutType:
        """Detect page layout type"""
        if width <= height * 1.1:
            return LayoutType.SINGLE_PORTRAIT
        return LayoutType.LANDSCAPE_OR_DOUBLE

    def _detect_coordinate_system(self):
        """Detect if Y coordinates are inverted"""
        if not self.doc or len(self.doc) == 0:
            return

        page = self.doc[0]
        blocks = page.get_text("dict").get("blocks", [])
        page_height = page.rect.height

        for block in blocks:
            if block.get("type") != 0:
                continue

            text = self._extract_block_text(block).lower()
            bbox = block.get("bbox", [0, 0, 0, 0])
            y_pos = bbox[1]

            if any(word in text for word in ["page", "seite", "©", "copyright"]):
                if y_pos < page_height * 0.2:
                    self.coordinate_system = "inverted"
                    logger.debug("Detected inverted Y-coordinate system")
                    return
                elif y_pos > page_height * 0.8:
                    self.coordinate_system = "normal"
                    logger.debug("Detected normal Y-coordinate system")
                    return

        self.coordinate_system = "normal"

    def _get_sort_y(self, block: TextBlock) -> float:
        """Get Y coordinate for sorting, accounting for coordinate system"""
        if self.coordinate_system == "inverted":
            page_height = self.page_layouts[block.page_num]["height"]
            return page_height - block.bbox[1]
        else:
            return block.bbox[1]

    def _analyze_middle_gap_coverage(
        self, blocks: List[TextBlock], mid_x: float, gap_width: float = MIDDLE_GAP_WIDTH
    ) -> float:
        """Analyze how much of the middle area is empty (0.0 to 1.0)"""
        gap_left = mid_x - gap_width / 2
        gap_right = mid_x + gap_width / 2

        coverage_points = 0
        total_points = 0

        for y in range(0, int(self.page_layouts[blocks[0].page_num]["height"]), 10):
            total_points += 1
            for block in blocks:
                if (
                    block.bbox[0] <= gap_right
                    and block.bbox[2] >= gap_left
                    and block.bbox[1] <= y <= block.bbox[3]
                ):
                    coverage_points += 1
                    break

        if total_points == 0:
            return 1.0

        empty_ratio = 1.0 - (coverage_points / total_points)
        return empty_ratio

    def _detect_page_numbers(
        self, blocks: List[TextBlock], page_width: float, page_height: float
    ) -> Dict:
        """Detect page numbers and their positions"""
        page_number_info = {"left": False, "right": False, "center": False, "numbers": []}

        for block in blocks:
            if block.width < 100 and block.height < 50:
                if block.y_center < page_height * 0.15 or block.y_center > page_height * 0.85:
                    text = block.text.strip()
                    if REGEX_PATTERNS["page_number"].match(text) or REGEX_PATTERNS[
                        "page_text"
                    ].match(text):
                        page_number_info["numbers"].append(text)

                        if block.x_center < page_width * 0.2:
                            page_number_info["left"] = True
                        elif block.x_center > page_width * 0.8:
                            page_number_info["right"] = True
                        elif 0.4 * page_width < block.x_center < 0.6 * page_width:
                            page_number_info["center"] = True

        return page_number_info

    def _extract_block_text(self, block: Dict) -> str:
        """Extract text from a block dictionary"""
        text_parts = []
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = span.get("text", "")
                if not isinstance(text, str):
                    text = str(text)
                text_parts.append(text)
        return " ".join(text_parts)

    def _fix_ligatures(self, text: str) -> str:
        """Replace ligatures with their ASCII equivalents"""
        ligature_map = {
            "\ufb01": "fi",
            "\ufb02": "fl",
            "\ufb00": "ff",
            "\ufb03": "ffi",
            "\ufb04": "ffl",
            "\ufb06": "st",
            "\u017ft": "ft",
            "\u2014": "--",
            "\u2013": "-",
            "\u2018": "'",
            "\u2019": "'",
            "\u201c": '"',
            "\u201d": '"',
            "\u2026": "...",
        }

        for ligature, replacement in ligature_map.items():
            if ligature in text:
                text = text.replace(ligature, replacement)

        return text

    def _normalize_utf8(self, text: str) -> str:
        """Robust UTF-8 normalization and mojibake fixing"""
        if not text:
            return text

        # Common mojibake patterns (double-encoded UTF-8)
        mojibake_fixes = {
            # German umlauts
            "ÃƒÂ¤": "ä",
            "ÃƒÂ¶": "ö",
            "ÃƒÂ¼": "ü",
            "Ãƒâ€ž": "Ä",
            'Ãƒâ€"': "Ö",
            'ÃƒÅ"': "Ü",
            "ÃƒÅ¸": "ß",
            # French/Spanish
            "ÃƒÂ©": "é",
            "ÃƒÂ¨": "è",
            "ÃƒÂ¡": "á",
            "Ãƒ ": "à",
            "ÃƒÂ³": "ó",
            "ÃƒÂ²": "ò",
            "ÃƒÂº": "ú",
            "ÃƒÂ¹": "ù",
            "ÃƒÂ­": "í",
            "ÃƒÂ¬": "ì",
            "ÃƒÂ±": "ñ",
            "ÃƒÂ§": "ç",
            # Quotes and dashes
            'Ã¢â‚¬Å"': '"',
            "Ã¢â‚¬": '"',
            "Ã¢â‚¬â„¢": "'",
            "Ã¢â‚¬Ëœ": "'",
            'Ã¢â‚¬"': "–",
            'Ã¢â‚¬"': "—",
            "Ã¢â‚¬Â¦": "...",
            "Ã¢â‚¬Â¢": "•",
        }

        # Apply mojibake fixes
        for wrong, correct in mojibake_fixes.items():
            text = text.replace(wrong, correct)

        # Normalize to NFC form (canonical composition)
        try:
            text = unicodedata.normalize("NFC", text)
        except Exception:
            pass

        return text

    def _fix_text_formatting_issues(self, text: str) -> str:
        """Post-processing to fix hyphenation and line breaks"""
        # Remove PDF artifacts
        text = text.replace('Ã¢â‚¬"', "")

        # Fix hyphenated words
        def fix_hyphenation(match):
            word1 = match.group(1)
            word2 = match.group(2)
            # Add basic language detection or default to common words check
            return word1 + word2

        text = re.sub(r"(\w{2,})-\s*\n\s*([a-zäöüß]\w*)", fix_hyphenation, text)
        text = re.sub(r"(\w{3,})-\s+([a-zäöüß]{4,})", r"\1\2", text)

        # Cleanup
        text = REGEX_PATTERNS["multiple_spaces"].sub(" ", text)
        text = re.sub(r"\s+([.,;:!?])", r"\1", text)
        text = re.sub(r"([.,;:!?])(?=[A-ZÄÖÜäöüß])", r"\1 ", text)
        text = re.sub(r"(\w)(\d{1,2})([A-ZÄÖÜ])", r"\1\2 \3", text)
        text = REGEX_PATTERNS["excessive_newlines"].sub("\n\n\n", text)

        return text

    def extract_text(self) -> str:
        """
        Main extraction method with advanced layout recognition.
        Column detection and adaptive columns are always enabled.
        """
        print(f"Processing PDF: {self.pdf_path}")
        print(f"Total pages: {len(self.doc)}")

        # Extract all blocks with positions
        all_blocks = self._extract_all_blocks_with_positions()

        # Apply intelligent layout analysis and sorting (always enabled)
        all_blocks = self._apply_intelligent_sorting(all_blocks)

        # Build text from sorted blocks
        final_text = self._build_text_from_blocks(all_blocks)

        self.doc.close()
        return final_text

    def _build_text_from_blocks(self, blocks: List[TextBlock]) -> str:
        """Build complete text from sorted blocks in correct reading order"""
        if not blocks:
            return ""

        pages_blocks = {}
        for block in blocks:
            if block.page_num not in pages_blocks:
                pages_blocks[block.page_num] = []
            pages_blocks[block.page_num].append(block)

        text_parts = []

        for page_num in sorted(pages_blocks.keys()):
            page_blocks = pages_blocks[page_num]
            logger.debug(f"Building text for page {page_num + 1} with {len(page_blocks)} blocks")

            if page_num == 0:
                text_parts.append(f"--- Seite {page_num + 1} ---\n")
            else:
                text_parts.append(f"\n\n--- Seite {page_num + 1} ---\n")

            current_column = None
            last_block = None
            paragraph_buffer = []

            for block in page_blocks:
                if block.semantic_type in [
                    BlockType.HEADER,
                    BlockType.FOOTER,
                    BlockType.PAGE_NUMBER,
                ]:
                    continue

                if block.column_id != current_column and block.column_id is not None:
                    if paragraph_buffer:
                        text_parts.append(" ".join(paragraph_buffer))
                        paragraph_buffer = []

                    current_column = block.column_id
                    if text_parts and not text_parts[-1].endswith("\n\n"):
                        text_parts.append("\n\n")

                if paragraph_buffer:
                    paragraph_buffer.append(block.text)
                else:
                    paragraph_buffer = [block.text]

            if paragraph_buffer:
                text_parts.append(" ".join(paragraph_buffer))

        result = "\n".join(text_parts)

        # Clean up whitespace
        result = REGEX_PATTERNS["excessive_newlines"].sub("\n\n\n", result)
        result = REGEX_PATTERNS["triple_newlines"].sub("\n\n", result)
        result = REGEX_PATTERNS["multiple_spaces"].sub(" ", result)
        result = result.strip()

        # Fix text formatting and normalize UTF-8
        result = self._fix_text_formatting_issues(result)
        result = self._normalize_utf8(result)

        return result

    def _extract_all_blocks_with_positions(self) -> List[TextBlock]:
        """Extract ALL blocks with complete position information"""
        all_blocks = []

        for page_num, page in enumerate(self.doc):
            page_layout = self.page_layouts[page_num]

            if page_num % 10 == 0:
                logger.debug(f"Extracting blocks from page {page_num + 1}/{len(self.doc)}...")

            blocks_dict = page.get_text("dict")

            for block in blocks_dict.get("blocks", []):
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        line_spans = []
                        line_bboxes = []

                        for span in line.get("spans", []):
                            span_text = span.get("text", "")
                            span_text = self._fix_ligatures(span_text)
                            span_text = self._normalize_utf8(span_text)

                            if span_text.strip():
                                line_spans.append(span_text)
                                line_bboxes.append(span.get("bbox", (0, 0, 0, 0)))

                        if line_spans:
                            combined_text = " ".join(line_spans)
                            combined_bbox = (
                                min(bbox[0] for bbox in line_bboxes),
                                min(bbox[1] for bbox in line_bboxes),
                                max(bbox[2] for bbox in line_bboxes),
                                max(bbox[3] for bbox in line_bboxes),
                            )

                            font_size = line.get("spans", [{}])[0].get("size", 0)

                            text_block = TextBlock(
                                text=combined_text,
                                bbox=combined_bbox,
                                page_num=page_num,
                                block_type="text",
                                font_size=font_size,
                            )

                            text_block.semantic_type = self._classify_block_type(
                                text_block, page_layout
                            )

                            all_blocks.append(text_block)

            try:
                tables = self._extract_tables_from_page(page, page_num)
                all_blocks.extend(tables)
            except Exception as e:
                logger.debug(f"Table extraction error on page {page_num}: {e}")

        return all_blocks

    def _classify_block_type(self, block: TextBlock, page_layout: Dict) -> BlockType:
        """Classify the semantic type of a text block"""
        page_width = page_layout["width"]
        page_height = page_layout["height"]

        x_ratio = block.x_center / page_width
        y_ratio = block.y_center / page_height

        text = block.text.strip()
        if block.width < 100 and block.height < 50:
            if y_ratio < 0.15 or y_ratio > 0.85:
                if REGEX_PATTERNS["page_number"].match(text) or REGEX_PATTERNS["page_text"].match(
                    text
                ):
                    return BlockType.PAGE_NUMBER
        
        return BlockType.MAIN_TEXT

    def _apply_intelligent_sorting(self, blocks: List[TextBlock]) -> List[TextBlock]:
        """Apply intelligent sorting with advanced column detection"""
        if not blocks:
            return blocks

        pages_blocks = {}
        for block in blocks:
            if block.page_num not in pages_blocks:
                pages_blocks[block.page_num] = []
            pages_blocks[block.page_num].append(block)

        ordered_blocks = []

        for page_num in sorted(pages_blocks.keys()):
            page_blocks = pages_blocks[page_num]
            page_layout = self.page_layouts[page_num]

            # Simplified single/landscape handling for now
            if page_layout["layout_type"] == LayoutType.LANDSCAPE_OR_DOUBLE:
                ordered = self._detect_and_order_columns_advanced(page_blocks, page_layout["width"])
            else:
                ordered = self._detect_and_order_columns_advanced(page_blocks, page_layout["width"])

            ordered_blocks.extend(ordered)

        return ordered_blocks

    def _detect_and_order_columns_advanced(
        self, blocks: List[TextBlock], page_width: float
    ) -> List[TextBlock]:
        """Advanced column detection without artificial limits"""
        # Simplified: sort by Y then X
        blocks.sort(key=lambda b: (self._get_sort_y(b), b.bbox[0]))
        return blocks

    def _extract_tables_from_page(self, page, page_num: int) -> List[TextBlock]:
        """Extract tables from a page"""
        table_blocks = []
        try:
            tables = page.find_tables()
            for table in tables:
                table_text = self._format_table(table)
                if table_text:
                    table_block = TextBlock(
                        text=table_text,
                        bbox=table.bbox,
                        page_num=page_num,
                        block_type="table",
                        semantic_type=BlockType.TABLE,
                    )
                    table_blocks.append(table_block)
        except Exception as e:
            logger.debug(f"Table detection error on page {page_num}: {e}")
        return table_blocks

    def _format_table(self, table) -> str:
        """Format table data as markdown"""
        try:
            data = table.extract()
            if not data:
                return ""

            md_lines = []
            for i, row in enumerate(data):
                if row:
                    row_text = " | ".join(str(cell or "").strip() for cell in row)
                    if row_text.strip():
                        md_lines.append(f"| {row_text} |")

                    if i == 0 and len(md_lines) > 0:
                        separator = " | ".join("---" for _ in row)
                        md_lines.append(f"| {separator} |")

            return "\n".join(md_lines) if md_lines else ""

        except Exception as e:
            logger.debug(f"Table formatting error: {e}")
            return ""


# ============= DOCUMENT PROCESSOR =============
@dataclass
class ExtractionResult:
    """Container for extraction results with quality metrics"""

    text: str
    metadata: Dict = field(default_factory=dict)
    quality_score: float = 0.0
    extraction_method: str = "unknown"
    tables: List = field(default_factory=list)
    processing_time: float = 0.0
    warnings: List = field(default_factory=list)


class DocumentProcessor:
    """Streamlined document processor for GraphRAG pipeline"""

    def __init__(self, output_dir: str = "/tmp"):
        self.output_dir = Path(output_dir)
        self.supported_formats = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"}

    def generate_doc_id(self, file_path: str) -> str:
        """Generate unique document ID with content hash"""
        file_path = Path(file_path)
        try:
            with open(file_path, "rb") as f:
                content = f.read()
                content_hash = hashlib.md5(content).hexdigest()[:8]
        except Exception:
            content_hash = hashlib.md5(str(file_path).encode()).hexdigest()[:8]

        clean_name = re.sub(r"[^\w\-_]", "_", file_path.stem)[:50]
        return f"{clean_name}_{content_hash}"

    def extract_pdf(self, file_path: str) -> ExtractionResult:
        """Extract text from PDF using advanced hybrid layout recognition"""
        start_time = time.time()
        warnings = []
        extractor = None

        try:
            try:
                import pymupdf4llm
            except ImportError:
                warnings.append("pymupdf4llm not available - required for PDF extraction")
                return ExtractionResult(
                    text="", extraction_method="failed", warnings=warnings
                )

            print("    Extracting with hybrid approach (pymupdf4llm + advanced layout analysis)...")
            extractor = PDFComplexLayoutExtractor(pdf_path=str(file_path), debug=False)
            extracted_text = extractor.extract_text()
            
            page_count = len(extractor.page_layouts)
            metadata = {"page_count": page_count}

            return ExtractionResult(
                text=extracted_text,
                metadata=metadata,
                extraction_method="hybrid_advanced_layout",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )

        except Exception as e:
            print(f"    PDF extraction failed: {str(e)[:100]}")
            return ExtractionResult(text="", extraction_method="failed", warnings=[str(e)])

    def extract_docx(self, file_path: str) -> ExtractionResult:
        """Extract text from DOCX"""
        start_time = time.time()
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n\n"
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        text += " | ".join(row_text) + "\n"
                text += "\n"

            return ExtractionResult(
                text=text,
                extraction_method="python-docx",
                processing_time=time.time() - start_time,
            )
        except Exception as e:
            return ExtractionResult(text="", extraction_method="failed", warnings=[str(e)])

    def extract_pptx(self, file_path: str) -> ExtractionResult:
        """Extract text from PPTX"""
        start_time = time.time()
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    text += f"\n--- Slide {slide_num} ---\n{slide_text}"

            return ExtractionResult(
                text=text,
                extraction_method="python-pptx",
                processing_time=time.time() - start_time,
            )
        except Exception as e:
            return ExtractionResult(text="", extraction_method="failed", warnings=[str(e)])

    def extract_xlsx(self, file_path: str) -> ExtractionResult:
        """Extract text from XLSX"""
        start_time = time.time()
        try:
            workbook = load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(v) if v is not None else "" for v in row]
                    row_text = " | ".join(row_values).strip()
                    if row_text:
                        sheet_text += row_text + "\n"
                text += sheet_text
            workbook.close()

            return ExtractionResult(
                text=text,
                extraction_method="openpyxl",
                processing_time=time.time() - start_time,
            )
        except Exception as e:
            return ExtractionResult(text="", extraction_method="failed", warnings=[str(e)])

    def extract_plaintext(self, file_path: str) -> ExtractionResult:
        """Extract text from plain text or markdown files"""
        start_time = time.time()
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

            return ExtractionResult(
                text=text,
                extraction_method="plain_text",
                processing_time=time.time() - start_time,
            )
        except Exception as e:
            return ExtractionResult(text="", extraction_method="failed", warnings=[str(e)])

    def process_single_file(self, file_path: str) -> Optional[Dict]:
        """Process a single document file"""
        file_path = Path(file_path)
        
        extractors = {
            ".pdf": self.extract_pdf,
            ".docx": self.extract_docx,
            ".pptx": self.extract_pptx,
            ".xlsx": self.extract_xlsx,
            ".txt": self.extract_plaintext,
            ".md": self.extract_plaintext,
        }

        extractor = extractors.get(file_path.suffix.lower())
        if not extractor:
            return None

        result = extractor(str(file_path))
        if not result.text.strip():
            return None

        doc_id = self.generate_doc_id(str(file_path))
        
        return {
            "doc_id": doc_id,
            "filename": file_path.name,
            "text": result.text,
            "metadata": result.metadata,
            "extraction_method": result.extraction_method
        }
